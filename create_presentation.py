"""Generate project presentation focused on model development."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent
OUTPUT = BASE / "Maize_Disease_Detection_Presentation.pptx"

# Scientific palette: white background, dark text, navy accents
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(33, 33, 33)
NAVY = RGBColor(0, 51, 102)
GRAY = RGBColor(80, 80, 80)
LIGHT_GRAY = RGBColor(230, 230, 230)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"


def set_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_title_bar(slide, title_text):
    """Add a simple top accent line and title."""
    line = slide.shapes.add_shape(
        1, Inches(0.5), Inches(0.35), Inches(9), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(9), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    return box


def add_bullets(slide, items, top=1.5, left=0.7, width=8.6, font_size=20, spacing=0.35):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = FONT_BODY
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.level = 0
        p.space_after = Pt(spacing * 12)
        p.line_spacing = 1.15


def add_table_slide(slide, headers, rows, top=1.6):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.8), Inches(top), Inches(8.4), Inches(0.45 * n_rows)
    )
    table = table_shape.table

    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_BODY
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    for i, row in enumerate(rows, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_BODY
                p.font.size = Pt(13)
                p.font.color.rgb = BLACK
                p.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: Title
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = NAVY
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(8.8), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Developing a Predictive Deep Learning Model for the Early Detection of Maize Leaf Diseases in Tanzania"
    p.font.name = FONT_TITLE
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    sub = tf.add_paragraph()
    sub.text = "Eastern Africa Statistical Training Centre"
    sub.font.name = FONT_BODY
    sub.font.size = Pt(18)
    sub.font.color.rgb = GRAY
    sub.alignment = PP_ALIGN.CENTER
    sub.space_before = Pt(24)

    team = tf.add_paragraph()
    team.text = (
        "Severia M. Banda | Moshi M. Kassim | Gidion F. Machumu | "
        "Rebeca P. Shauri | Edgar S. Materu"
    )
    team.font.name = FONT_BODY
    team.font.size = Pt(14)
    team.font.color.rgb = GRAY
    team.alignment = PP_ALIGN.CENTER
    team.space_before = Pt(18)

    meta = tf.add_paragraph()
    meta.text = "Bachelor of Data Science, Year III | Academic Year 2025/2026"
    meta.font.name = FONT_BODY
    meta.font.size = Pt(14)
    meta.font.color.rgb = GRAY
    meta.alignment = PP_ALIGN.CENTER
    meta.space_before = Pt(8)

    sup = tf.add_paragraph()
    sup.text = "Supervisor: Mr Kelvin Rushobola"
    sup.font.name = FONT_BODY
    sup.font.size = Pt(14)
    sup.font.color.rgb = GRAY
    sup.alignment = PP_ALIGN.CENTER
    sup.space_before = Pt(8)

    # Slide 2: Background
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Background")
    add_bullets(
        slide,
        [
            "Maize is the primary staple crop in Tanzania and is highly vulnerable to foliar diseases.",
            "Common Rust, Northern Leaf Blight, and Cercospora Leaf Spot can cause substantial yield losses if not detected early.",
            "Current diagnosis depends on manual field scouting and extension officers, which is slow and limited by staff shortages.",
            "Convolutional Neural Networks can learn disease patterns from labelled leaf images and support automated classification.",
            "This project develops and evaluates such a model for maize leaf disease detection.",
        ],
        font_size=19,
    )

    # Slide 3: General Objective
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "General Objective")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "To develop a deep learning image classification model capable of automatically "
        "identifying and diagnosing common maize leaf diseases, and to deploy the model "
        "in an interactive application for real-time diagnosis."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(22)
    p.font.color.rgb = BLACK
    p.line_spacing = 1.3

    # Slide 4: Specific Objectives
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Specific Objectives")
    add_bullets(
        slide,
        [
            "Source a standardized dataset of healthy and diseased maize leaf images.",
            "Train Convolutional Neural Network models for classifying distinct maize foliar diseases.",
            "Evaluate predictive performance using accuracy, precision, recall, F1-score, and confusion matrix.",
            "Deploy the selected model in an interactive graphical interface for real-time image diagnosis.",
        ],
        font_size=20,
    )

    # Slide 5: Research Questions
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Research Questions")
    add_bullets(
        slide,
        [
            "What image preprocessing techniques are required to prepare raw leaf data for deep learning?",
            "How accurately can a Convolutional Neural Network classify different maize leaf disease classes?",
            "How can the predictive model be deployed to accept new images and output diagnostic confidence scores?",
        ],
        font_size=20,
        spacing=0.5,
    )

    # Slide 6: Methodology Overview
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Methodology Overview")
    add_bullets(
        slide,
        [
            "Approach: Predictive analytics and computer vision using supervised image classification.",
            "Data source: PlantVillage maize subset (4,188 labelled RGB images).",
            "Development environment: Python 3.10+, TensorFlow/Keras, OpenCV, scikit-learn.",
            "Workflow: Data acquisition, preprocessing, model training, evaluation, model selection, deployment.",
            "Two architectures were compared: a Custom CNN and ResNet50 transfer learning.",
        ],
        font_size=19,
    )

    # Slide 7: Dataset
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Dataset and Sampling")
    add_bullets(
        slide,
        [
            "Four disease classes: Healthy, Common Rust, Northern Leaf Blight, Cercospora Leaf Spot.",
            "Total sample: 4,188 images from the PlantVillage repository.",
            "Stratified random split to preserve class proportions across sets.",
        ],
        font_size=18,
        top=1.45,
    )
    add_table_slide(
        slide,
        ["Split", "Images", "Percentage"],
        [
            ["Training", "2,931", "70%"],
            ["Validation", "628", "15%"],
            ["Test", "629", "15%"],
            ["Total", "4,188", "100%"],
        ],
        top=3.2,
    )

    # Slide 8: Preprocessing
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Data Preprocessing")
    add_bullets(
        slide,
        [
            "Resizing: All images resized to 128 x 128 pixels.",
            "Colour conversion: RGB converted to BGR to match the training pipeline.",
            "Normalization: Pixel values scaled from 0 to 255 down to the range 0 to 1.",
            "Batch formatting: Input tensor shape (1, 128, 128, 3) for model inference.",
            "Data augmentation was applied during training to reduce overfitting.",
        ],
        font_size=19,
    )

    # Slide 9: Custom CNN Architecture
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Model 1: Custom CNN Architecture")
    add_bullets(
        slide,
        [
            "Input layer: 128 x 128 x 3",
            "Block 1: Conv2D(32) + Conv2D(32) + MaxPooling + Dropout",
            "Block 2: Conv2D(64) + Conv2D(64) + MaxPooling + Dropout",
            "Block 3: Conv2D(128) + MaxPooling + Dropout",
            "Classifier: Flatten + Dense(256) + Dropout + Dense(4, softmax)",
            "Optimizer: Adam | Loss: Sparse Categorical Cross-Entropy",
            "Training: 25 epochs, batch size 32",
        ],
        font_size=17,
        spacing=0.22,
    )

    # Slide 10: ResNet50 Baseline
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Model 2: ResNet50 Baseline")
    add_bullets(
        slide,
        [
            "Transfer learning using ResNet50 pre-trained on ImageNet.",
            "Base layers frozen; only the classification head was trained.",
            "Classification head: GlobalAveragePooling + Dense(256) + Dense(128) + Dense(4).",
            "Optimizer: Adam with learning rate 0.0001.",
            "Training: 100 epochs on the same dataset split.",
            "Purpose: Benchmark comparison against the task-specific Custom CNN.",
        ],
        font_size=19,
    )

    # Slide 11: Training Process
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Model Training and Validation")
    add_bullets(
        slide,
        [
            "Training performed on 2,931 images with validation monitoring during each epoch.",
            "Validation set (628 images) used to track loss and accuracy during training.",
            "Final evaluation conducted on a held-out test set of 629 images.",
            "Test images were not used during training or hyperparameter decisions.",
            "Metrics recorded: accuracy, loss, precision, recall, F1-score, and confusion matrix.",
            "Best-performing model saved as corn_disease_cnn.h5 for deployment.",
        ],
        font_size=18,
    )

    # Slide 12: Model Comparison
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Model Comparison and Selection")
    add_table_slide(
        slide,
        ["Model", "Test Accuracy", "Test Loss", "Selected"],
        [
            ["Custom CNN", "92.37%", "0.3725", "Yes"],
            ["ResNet50 (frozen)", "77.42%", "0.5277", "No"],
        ],
        top=2.0,
    )
    note = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(8.4), Inches(2.5))
    tf = note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "The Custom CNN outperformed the frozen ResNet50 baseline on the same test set. "
        "For this four-class problem with approximately 4,000 images, a task-specific "
        "architecture learned maize disease features more effectively than a general "
        "pre-trained network with frozen layers."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK
    p.line_spacing = 1.3

    # Slide 13: Findings - Overall Performance
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Findings: Overall Model Performance")
    add_bullets(
        slide,
        [
            "Custom CNN achieved 92.37% test accuracy on 629 unseen images.",
            "Healthy and Common Rust classes were classified with near-perfect performance.",
            "Northern Leaf Blight showed strong recall (0.94) with moderate precision (0.83).",
            "Cercospora Leaf Spot was the most challenging class (recall 0.65).",
            "Results confirm that deep learning can classify maize foliar diseases from digital images.",
        ],
        font_size=18,
    )

    # Slide 14: Per-class metrics
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Findings: Per-Class Performance (Custom CNN)")
    add_table_slide(
        slide,
        ["Class", "Precision", "Recall", "F1-Score"],
        [
            ["Common Rust", "0.97", "0.96", "0.97"],
            ["Northern Leaf Blight", "0.83", "0.94", "0.88"],
            ["Healthy", "1.00", "1.00", "1.00"],
            ["Cercospora Leaf Spot", "0.85", "0.65", "0.74"],
        ],
        top=1.7,
    )
    note = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.4), Inches(2.0))
    tf = note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "The main classification error was confusion between Cercospora Leaf Spot and "
        "Northern Leaf Blight, consistent with visual similarity in lesion appearance."
    )
    p.font.name = FONT_BODY
    p.font.size = Pt(17)
    p.font.color.rgb = BLACK

    # Slide 15: Confusion Matrix
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Findings: Confusion Matrix (Test Set)")
    add_table_slide(
        slide,
        ["Actual \\ Predicted", "NLB", "Rust", "CLS", "Healthy"],
        [
            ["Northern Leaf Blight", "162", "5", "5", "0"],
            ["Common Rust", "3", "188", "5", "0"],
            ["Cercospora Leaf Spot", "22", "8", "56", "0"],
            ["Healthy", "0", "0", "0", "175"],
        ],
        top=1.65,
    )

    # Slide 16: Discussion
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Discussion")
    add_bullets(
        slide,
        [
            "Results align with prior plant disease studies using deep learning on PlantVillage data.",
            "A simpler Custom CNN exceeded a frozen ResNet50, suggesting domain-specific architectures can be effective at this dataset scale.",
            "Strong performance on Healthy and Common Rust supports reliable detection of clear symptom patterns.",
            "Lower recall for Cercospora Leaf Spot indicates need for additional training samples and augmentation.",
            "Laboratory accuracy (92.37%) is competitive, though field validation on Tanzanian farm images remains necessary.",
        ],
        font_size=17,
    )

    # Slide 17: Limitations
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Limitations")
    add_bullets(
        slide,
        [
            "Dataset images originate from controlled PlantVillage conditions, not Tanzanian field environments.",
            "Input resolution of 128 x 128 pixels may reduce fine lesion detail.",
            "Cercospora Leaf Spot classification requires further improvement.",
            "Model performance under varied lighting, backgrounds, and camera quality has not yet been validated in the field.",
        ],
        font_size=19,
    )

    # Slide 18: Conclusion
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    add_title_bar(slide, "Conclusion")
    add_bullets(
        slide,
        [
            "A standardized dataset of 4,188 maize leaf images was successfully sourced and prepared.",
            "Custom CNN and ResNet50 models were trained; the Custom CNN achieved 92.37% test accuracy.",
            "Model performance was rigorously evaluated using standard classification metrics.",
            "The selected model was saved and integrated for real-time disease classification.",
            "The project demonstrates that a task-specific CNN can provide an accessible diagnostic support tool for maize foliar diseases in Tanzania.",
        ],
        font_size=17,
    )

    # Slide 19: Thank you
    slide = prs.slides.add_slide(blank)
    set_white_background(slide)
    accent = slide.shapes.add_shape(1, Inches(0), Inches(7.38), Inches(10), Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = NAVY
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(8.8), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.name = FONT_TITLE
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    q = tf.add_paragraph()
    q.text = "Questions and Discussion"
    q.font.name = FONT_BODY
    q.font.size = Pt(20)
    q.font.color.rgb = GRAY
    q.alignment = PP_ALIGN.CENTER
    q.space_before = Pt(16)

    prs.save(OUTPUT)
    print(f"Presentation saved: {OUTPUT}")


if __name__ == "__main__":
    build_presentation()
